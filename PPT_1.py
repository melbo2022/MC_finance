from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import streamlit as st

# Excelファイルとシート名を指定
excel_file = 'DC_data.xlsx'
sheet_name = 'Table_1'

# Excelファイルからデータを読み込む
df = pd.read_excel(excel_file, sheet_name=sheet_name)

# A1:D10の範囲を抽出
df = df.iloc[0:5, 0:4]

# データを表示
st.write(df)

# 新しいパワーポイントプレゼンテーションを作成
prs = Presentation()

# 新しいスライドを追加
slide = prs.slides.add_slide(prs.slide_layouts[5])

# スライドにテーブルを追加
table = slide.shapes.add_table(df.shape[0], df.shape[1], Inches(2.0), Inches(2.0), Inches(6.0), Inches(0.8)).table

# テーブルにデータを追加
for i in range(df.shape[0]):
    for j in range(df.shape[1]):
        table.cell(i, j).text = str(df.iat[i, j])

# パワーポイントプレゼンテーションを保存
prs.save('output.pptx')


